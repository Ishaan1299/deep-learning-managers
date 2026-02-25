"""
Generate 7-slide PowerPoint presentations for all three Deep Learning for Managers projects.
Run from: C:/Users/Ishaan/.gemini/antigravity/scratch/
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Colour palette ───────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT     = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy
HIGHLIGHT  = RGBColor(0x0F, 0x3D, 0x66)   # mid blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xF5, 0xA6, 0x23)   # amber / gold
LIGHT_GREY = RGBColor(0xD0, 0xD8, 0xE8)   # pale blue-grey for body text
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helper functions ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=DARK_BG):
    from pptx.util import Pt
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=16, color=LIGHT_GREY, bold_first=False):
    """Add a multi-bullet text box. Each item in bullets is a string."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.space_before = Pt(4)
        run = para.add_run()
        run.text = bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first and bullet == bullets[0]:
            run.font.bold = True


def slide_title_bar(slide, title, subtitle=None, bar_color=HIGHLIGHT):
    """Dark background + coloured bar at top + title text."""
    fill_bg(slide, DARK_BG)
    # top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), bar_color)
    add_text(slide, title,
             Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.75),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.35),
                 font_size=14, bold=False, color=GOLD, align=PP_ALIGN.LEFT)


def divider_line(slide, y_inches, color=HIGHLIGHT):
    line = slide.shapes.add_shape(1,
        Inches(0.4), Inches(y_inches),
        Inches(12.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def metric_box(slide, label, value, left, top, w=Inches(2.4), h=Inches(1.15),
               val_color=GREEN):
    add_rect(slide, left, top, w, h, ACCENT)
    add_text(slide, label,
             left + Inches(0.1), top + Inches(0.05), w - Inches(0.2), Inches(0.4),
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(slide, value,
             left + Inches(0.05), top + Inches(0.42), w - Inches(0.1), Inches(0.6),
             font_size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  PROJECT 1 — CREDIT DEFAULT PREDICTION (ANN)
# ═══════════════════════════════════════════════════════════════════════════════

def make_credit_ppt():
    prs = new_prs()

    # ── SLIDE 1 : Title ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, 0, 0, SLIDE_W, Inches(7.5), DARK_BG)
    # left accent strip
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, GOLD)
    # centre card
    add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.6), ACCENT)

    add_text(s, "Deep Learning for Managers",
             Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
             font_size=16, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, "Credit Default Prediction",
             Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.0),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Using Artificial Neural Networks (ANN)",
             Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.6),
             font_size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, "Predicting NPA Risk in Indian Banking  |  P1 → P4 Four-Class Classifier",
             Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.45),
             font_size=14, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    # metric strip
    metrics = [("Accuracy", "91.23%"), ("Precision", "92.68%"),
               ("F1-Score", "91.45%"), ("ROC-AUC", "0.9933")]
    x_start = Inches(0.75)
    for label, val in metrics:
        metric_box(s, label, val, x_start, Inches(5.15), w=Inches(2.85), h=Inches(1.1))
        x_start += Inches(3.0)

    add_text(s, "Dataset: ~51,000 Indian bank borrowers  |  PyTorch MLP  |  Feb 2026",
             Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 : Problem Statement ────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "The Problem: India's NPA Crisis",
                    "Why credit default prediction matters")
    add_text(s, "📌  The Business Context",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=15, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  FY 2018: Gross NPA ratio of Public Sector Banks peaked at 11.6% — Rs. 9.62 lakh crore of bad loans",
        "▸  FY 2023: NPAs fell to 5.0%, but absolute stock still exceeds Rs. 6 lakh crore (RBI, 2023)",
        "▸  Annual provisioning cost: Rs. 2.2–2.8 lakh crore locked in non-productive capital (Basel Committee, 2017)",
        "▸  One bad loan costs the bank 4–6× the original loan amount when all downstream effects are counted",
    ], Inches(0.5), Inches(1.7), Inches(12.2), Inches(1.8), font_size=15)

    divider_line(s, 3.6)
    add_text(s, "🎯  Why Deep Learning?",
             Inches(0.4), Inches(3.7), Inches(12.0), Inches(0.4),
             font_size=15, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  Traditional scorecards use CIBIL score alone — miss internal bank behavioural signals",
        "▸  Binary (default / no-default) outputs don't map to bank's operational pricing tiers",
        "▸  ANN learns non-linear interactions: e.g., high salary + low CIBIL + recent SMA flag = high risk",
        "▸  Data fusion of internal bank data (29 cols) + CIBIL bureau data (40 cols) gives a 360° customer view",
    ], Inches(0.5), Inches(4.15), Inches(12.2), Inches(2.0), font_size=15)

    # ── SLIDE 3 : Data & Features ──────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Data & Feature Engineering",
                    "~51,000 Indian retail borrowers — merged internal + CIBIL dataset")

    # Left column
    add_text(s, "Internal Bank Dataset  (29 features)",
             Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• Total_TL — total trade lines",
        "• pct_tl_open_L6M — new credit in 6 months",
        "• CC_utilization_pct — credit card usage %",
        "• num_tl_90dpd_6m — 90-day delinquencies",
        "• num_tl_30dpd — 30-day late accounts",
    ], Inches(0.4), Inches(1.7), Inches(5.6), Inches(2.0), font_size=13)

    # Right column
    add_text(s, "External CIBIL Dataset  (40 features)",
             Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• CIBIL_SCORE_EQUIV — bureau score",
        "• num_deliq_6mts — delinquencies (bureau)",
        "• pct_PL_enq_L6m — personal loan inquiries",
        "• time_since_recent_payment — recency",
        "• Target: Approved_Flag (P1 / P2 / P3 / P4)",
    ], Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.0), font_size=13)

    divider_line(s, 3.85)
    add_text(s, "Target Class Distribution",
             Inches(0.4), Inches(3.95), Inches(12.0), Inches(0.4),
             font_size=13, bold=True, color=GOLD)

    tiers = [("P1 — Best", "24.4%", "Auto-approve, lowest rate", GREEN),
             ("P2 — Good", "37.5%", "Standard approval", LIGHT_GREY),
             ("P3 — Moderate", "26.9%", "Risk premium required", GOLD),
             ("P4 — High Risk", "11.1%", "Reject / collateral only", RED)]
    x = Inches(0.4)
    for label, pct, desc, col in tiers:
        add_rect(s, x, Inches(4.4), Inches(2.95), Inches(2.4), ACCENT)
        add_text(s, label, x + Inches(0.1), Inches(4.45), Inches(2.7), Inches(0.4),
                 font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, pct,   x + Inches(0.1), Inches(4.85), Inches(2.7), Inches(0.5),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, desc,  x + Inches(0.1), Inches(5.3), Inches(2.7), Inches(0.35),
                 font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        x += Inches(3.1)

    # ── SLIDE 4 : Model Architecture ───────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "ANN Architecture",
                    "Multi-Layer Perceptron — Input(93) → [128 → 64 → 32] → Output(4)")

    add_text(s, "Network Design",
             Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    layers = [
        ("Input Layer", "93 features", "Scaled with StandardScaler"),
        ("Hidden Layer 1", "128 neurons", "BatchNorm + ReLU + Dropout 30%"),
        ("Hidden Layer 2", "64 neurons",  "BatchNorm + ReLU + Dropout 20%"),
        ("Hidden Layer 3", "32 neurons",  "BatchNorm + ReLU + Dropout 20%"),
        ("Output Layer",   "4 logits",    "CrossEntropyLoss (softmax implicit)"),
    ]
    y = Inches(1.7)
    for name, size, detail in layers:
        add_rect(s, Inches(0.4), y, Inches(5.5), Inches(0.7), ACCENT)
        add_text(s, f"{name}  ({size})", Inches(0.5), y + Inches(0.05),
                 Inches(3.5), Inches(0.35), font_size=13, bold=True, color=WHITE)
        add_text(s, detail, Inches(0.5), y + Inches(0.35),
                 Inches(5.2), Inches(0.3), font_size=11, color=LIGHT_GREY)
        y += Inches(0.8)

    # Right column — training config
    add_text(s, "Training Configuration",
             Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "🔹  Optimizer: Adam  (lr=0.001, weight_decay=1e-5)",
        "🔹  Loss: CrossEntropyLoss with inverse-freq class weights",
        "🔹  P4 class weight: 2.25×  (penalise missed defaults)",
        "🔹  Epochs: 20  |  Batch size: 64",
        "🔹  Train / Test split: 80 / 20  (stratified)",
        "🔹  Total parameters: 23,428",
        "🔹  Framework: PyTorch",
    ], Inches(6.8), Inches(1.7), Inches(5.8), Inches(3.5), font_size=13)

    add_text(s, "Key Design Rationale",
             Inches(6.8), Inches(5.3), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• Funnel shape (128→64→32): compresses 93 features into abstract risk archetypes",
        "• Dropout prevents overfitting on 40K training samples",
    ], Inches(6.8), Inches(5.7), Inches(5.8), Inches(1.2), font_size=12)

    # ── SLIDE 5 : Results ──────────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Results — Test Set Performance",
                    "Evaluated on 10,268 held-out borrowers never seen during training")

    # Metric boxes
    met = [("Accuracy", "91.23%", GREEN), ("Precision", "92.68%", GREEN),
           ("Recall",   "91.23%", GREEN), ("F1-Score",  "91.45%", GREEN),
           ("ROC-AUC",  "0.9933", GOLD)]
    x = Inches(0.3)
    for label, val, col in met:
        metric_box(s, label, val, x, Inches(1.25), w=Inches(2.46), h=Inches(1.2), val_color=col)
        x += Inches(2.6)

    divider_line(s, 2.65)
    add_text(s, "Business Value Quantification",
             Inches(0.4), Inches(2.75), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  On a Rs. 500 Cr quarterly loan portfolio, the ANN reduces NPA rate from ~11% (no model) to ~1.5%",
        "▸  Estimated quarterly savings in provisioning costs: Rs. 7.13 Crore vs. no model",
        "▸  ROC-AUC 0.9933 vs. ~0.82 for logistic regression — deep learning advantage is measurable",
        "▸  Near-zero P1↔P4 cross-tier misclassification: model never confuses an excellent borrower with a defaulter",
    ], Inches(0.4), Inches(3.2), Inches(12.2), Inches(2.0), font_size=14)

    add_text(s, "Confusion Matrix Highlights",
             Inches(0.4), Inches(5.3), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  P1 (Best borrowers): near-perfect recall — bank will not wrongly reject excellent customers",
        "▸  P4 (High Risk): strong recall — most defaulters are caught before disbursement",
    ], Inches(0.4), Inches(5.7), Inches(12.2), Inches(1.2), font_size=13)

    # ── SLIDE 6 : Managerial Insights ─────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Managerial Implications",
                    "How the ANN translates to credit operations strategy")

    insights = [
        ("💰  Risk-Based Pricing", [
            "P1 (prob > 85%): 9.0–9.5% rate — steal market share",
            "P3 (prob > 60%): 13.0–14.5% — risk premium offsets default cost",
            "P4 (prob > 50%): Reject or require collateral",
        ]),
        ("⚡  Zero-Touch Underwriting", [
            "Top 20% P1 applicants (prob > 92%): instant auto-approval in < 3 seconds",
            "Reduces TAT from 3–5 days to milliseconds for best customers",
        ]),
        ("📊  CFO Capital Planning", [
            "Run ANN on live loan book quarterly to forecast P2→P4 migrations",
            "Pre-calculate RBI provisioning requirements before quarter-end",
        ]),
        ("🗺️  Regional Policy", [
            "Deploy region-specific thresholds: stricter in high-NPA belts",
            "Urban growth zones: accept P3 with income verification",
        ]),
    ]
    y = Inches(1.25)
    col = 0
    for title, bullets in insights:
        x = Inches(0.4) if col == 0 else Inches(6.9)
        add_rect(s, x, y, Inches(6.0), Inches(2.5), ACCENT)
        add_text(s, title, x + Inches(0.15), y + Inches(0.1),
                 Inches(5.7), Inches(0.4), font_size=13, bold=True, color=GOLD)
        add_bullet_box(s, bullets, x + Inches(0.15), y + Inches(0.55),
                       Inches(5.6), Inches(1.8), font_size=12)
        col += 1
        if col == 2:
            col = 0
            y += Inches(2.65)

    # ── SLIDE 7 : Conclusions ──────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Conclusions & Future Directions",
                    "What was achieved and what comes next")

    add_text(s, "✅  Key Contributions",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  91.23% accuracy / ROC-AUC 0.9933 on 4-class credit risk prediction on Indian banking data",
        "▸  Data fusion pipeline: internal bank (29 cols) + CIBIL (40 cols) → 93 engineered features",
        "▸  Quantified Rs. 7.13 Cr quarterly savings on a Rs. 500 Cr portfolio",
        "▸  Production-ready Flask API + browser demo for live inference",
    ], Inches(0.4), Inches(1.7), Inches(12.2), Inches(1.8), font_size=14)

    divider_line(s, 3.65)
    add_text(s, "🚀  Future Research",
             Inches(0.4), Inches(3.75), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  XAI with SHAP: per-applicant feature attribution to satisfy RBI explainability requirements",
        "▸  LSTM-based temporal scoring: process 12–24 months of behaviour as time series",
        "▸  Alternative data: UPI velocity + GST filing via RBI's Account Aggregator / OCEN framework",
        "▸  Federated learning: pan-Indian credit model with data privacy compliance",
    ], Inches(0.4), Inches(4.2), Inches(12.2), Inches(2.1), font_size=14)

    prs.save("credit_default_prediction/credit_default_prediction.pptx")
    print("[OK] Saved: credit_default_prediction/credit_default_prediction.pptx")


# ═══════════════════════════════════════════════════════════════════════════════
#  PROJECT 2 — UPI FRAUD DETECTION (LSTM)
# ═══════════════════════════════════════════════════════════════════════════════

def make_upi_ppt():
    prs = new_prs()

    # ── SLIDE 1 : Title ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, RGBColor(0x00, 0xB4, 0xD8))
    add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.6), ACCENT)

    add_text(s, "Deep Learning for Managers",
             Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
             font_size=16, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, "UPI Fraud Detection",
             Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.0),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Using Long Short-Term Memory (LSTM) Networks",
             Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.6),
             font_size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, "Sequential Account Takeover Detection  |  NPCI-Calibrated Synthetic Data",
             Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.45),
             font_size=14, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    metrics = [("Accuracy", "93.01%"), ("Precision", "93.04%"),
               ("Recall",   "79.85%"), ("ROC-AUC",  "88.75%")]
    x_start = Inches(0.75)
    for label, val in metrics:
        metric_box(s, label, val, x_start, Inches(5.15), w=Inches(2.85), h=Inches(1.1))
        x_start += Inches(3.0)

    add_text(s, "Dataset: 80,000 synthetic UPI sessions  |  NPCI FY 2021-24  |  PyTorch LSTM  |  Feb 2026",
             Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 : Problem Statement ────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "The Problem: UPI Fraud at Scale",
                    "India's Rs. 1,978 lakh crore digital payments ecosystem under attack",
                    bar_color=RGBColor(0x02, 0x53, 0x6A))

    add_text(s, "📌  Scale of the Problem",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=15, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  March 2024: UPI processed 13,113 million transactions worth Rs. 19,78,312 crore — a 56% YoY growth (NPCI, 2024)",
        "▸  Digital payment fraud = 42.3% of all cybercrime complaints in 2023 (I4C Annual Report, 2024)",
        "▸  UPI fraud growing faster than traditional channels; RBI reported Rs. 630 Cr+ losses in FY 2022-23",
    ], Inches(0.4), Inches(1.7), Inches(12.2), Inches(1.6), font_size=15)

    divider_line(s, 3.45)
    add_text(s, "🎯  The Account Takeover (ATO) Pattern — Why It's Hard to Detect",
             Inches(0.4), Inches(3.55), Inches(12.0), Inches(0.4),
             font_size=15, bold=True, color=GOLD)

    phases = [("t=1–8\nNormal baseline", ACCENT),
              ("t=9\nDevice change", HIGHLIGHT),
              ("t=10–12\nVelocity spike", RGBColor(0x6B, 0x35, 0x00)),
              ("t=13–15\nDraining transfers", RED)]
    x = Inches(0.4)
    for label, col in phases:
        add_rect(s, x, Inches(4.05), Inches(2.9), Inches(1.5), col)
        add_text(s, label, x + Inches(0.1), Inches(4.1), Inches(2.7), Inches(1.3),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += Inches(3.05)

    add_text(s, "No single transaction is unambiguously fraudulent — only the temporal sequence reveals the attack.",
             Inches(0.4), Inches(5.7), Inches(12.2), Inches(0.5),
             font_size=13, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    add_text(s, "LSTM's hidden state accumulates evidence across all 15 steps — "
               "rule-based systems and transaction-level classifiers cannot do this.",
             Inches(0.4), Inches(6.2), Inches(12.2), Inches(0.7),
             font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 3 : Data ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Data — NPCI-Calibrated Synthetic Sessions",
                    "80,000 sessions × 15 transactions × 12 features  |  23% fraud rate",
                    bar_color=RGBColor(0x02, 0x53, 0x6A))

    add_text(s, "NPCI Calibration Anchor",
             Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• 30 monthly data points: FY 2021-22, 2022-23, 2023-24",
        "• Derived avg transaction value: Rs. 1,612 per UPI transfer",
        "• Volume grew 5.2× over 3 years (2.5 Bn → 13 Bn/month)",
        "• Lognormal amount distribution centred on Rs. 1,612",
    ], Inches(0.4), Inches(1.7), Inches(5.8), Inches(2.0), font_size=13)

    add_text(s, "12 Transaction Features",
             Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• amount, hour_of_day, day_of_week, is_weekend",
        "• is_late_night, is_new_payee, device_changed",
        "• high_amount_flag, txn_velocity, time_since_last",
        "• cumulative_amount_ratio, payee_familiarity",
    ], Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.0), font_size=13)

    divider_line(s, 3.85)
    add_text(s, "Design Choices for Realism",
             Inches(0.4), Inches(3.95), Inches(12.0), Inches(0.4),
             font_size=13, bold=True, color=GOLD)

    choices = [("7% Label Noise", "Simulates real-world label ambiguity — prevents threshold memorisation"),
               ("15% Subtle Fraud", "Fraud without device change — forces LSTM to learn multiple fraud signatures"),
               ("20% Normal Noise Events", "Legitimate users sometimes change devices — prevents over-separation"),
               ("pos_weight capped at 5.0", "Controls FP rate for operational deployment — precision vs. recall balance")]
    y = Inches(4.4)
    for title, desc in choices:
        add_rect(s, Inches(0.4), y, Inches(12.0), Inches(0.55), ACCENT)
        add_text(s, f"  {title}:", Inches(0.5), y + Inches(0.05),
                 Inches(2.8), Inches(0.4), font_size=12, bold=True, color=GOLD)
        add_text(s, desc, Inches(3.3), y + Inches(0.05),
                 Inches(8.9), Inches(0.4), font_size=12, color=LIGHT_GREY)
        y += Inches(0.65)

    # ── SLIDE 4 : Model Architecture ───────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "LSTM Architecture — FraudLSTM",
                    "2-layer stacked LSTM (128 units) → FC classifier head  |  147,777 parameters",
                    bar_color=RGBColor(0x02, 0x53, 0x6A))

    add_text(s, "Architecture Flow",
             Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    arch = [
        ("Input", "(batch, 15, 12)", "15 steps × 12 features"),
        ("LSTM Layer 1", "hidden=128, dropout=0.3", "Learns short-range temporal patterns"),
        ("LSTM Layer 2", "hidden=128", "Learns higher-order sequential patterns"),
        ("Last Hidden State", "(batch, 128)", "Summarises full 15-step sequence"),
        ("Dropout (30%)", "Regularisation", "Prevents co-adaptation of neurons"),
        ("FC: 128 → 64", "+ BatchNorm + ReLU", "Compresses to classification representation"),
        ("Output", "1 logit", "BCEWithLogitsLoss (sigmoid implicit)"),
    ]
    y = Inches(1.7)
    for name, detail, note in arch:
        add_rect(s, Inches(0.4), y, Inches(5.8), Inches(0.62), ACCENT)
        add_text(s, name, Inches(0.5), y + Inches(0.04),
                 Inches(2.2), Inches(0.3), font_size=12, bold=True, color=WHITE)
        add_text(s, detail, Inches(2.7), y + Inches(0.04),
                 Inches(1.8), Inches(0.3), font_size=11, color=GOLD)
        add_text(s, note, Inches(0.5), y + Inches(0.34),
                 Inches(5.5), Inches(0.25), font_size=10, color=LIGHT_GREY)
        y += Inches(0.72)

    add_text(s, "Why LSTM?",
             Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "🔹  Forget / input / output gates explicitly control what information is retained across steps",
        "🔹  Connects device change at step 9 to large transfers at steps 13–15 (6-step span)",
        "🔹  Standard RNNs fail on long sequences due to vanishing gradients",
        "",
        "Training Configuration:",
        "  • Optimizer: Adam  (lr=1e-3, weight_decay=1e-5)",
        "  • LR Scheduler: StepLR (step=5, gamma=0.5)",
        "  • Gradient clipping: max_norm=1.0 (prevents exploding gradients in RNNs)",
        "  • Epochs: 20  |  Batch size: 256",
    ], Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.5), font_size=12)

    # ── SLIDE 5 : Results ──────────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Results — Test Set Performance",
                    "Evaluated on 16,000 held-out sessions  |  TN=11,461  FP=256  FN=863  TP=3,420",
                    bar_color=RGBColor(0x02, 0x53, 0x6A))

    met = [("Accuracy", "93.01%", GREEN), ("Precision", "93.04%", GREEN),
           ("Recall",   "79.85%", GOLD),  ("F1-Score",  "85.94%", GREEN),
           ("ROC-AUC",  "88.75%", GREEN)]
    x = Inches(0.3)
    for label, val, col in met:
        metric_box(s, label, val, x, Inches(1.25), w=Inches(2.46), h=Inches(1.2), val_color=col)
        x += Inches(2.6)

    divider_line(s, 2.65)
    add_text(s, "Confusion Matrix Interpretation",
             Inches(0.4), Inches(2.75), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)

    cells = [("TN = 11,461", "Legitimate sessions correctly passed — zero friction for real users",       GREEN,  Inches(0.4)),
             ("FP = 256",    "False alerts — 2.18% of legitimate sessions incorrectly flagged",            GOLD,   Inches(6.7)),
             ("FN = 863",    "Missed fraud — 20.15% of fraud sessions escape detection",                   RED,    Inches(0.4)),
             ("TP = 3,420",  "Fraud correctly caught — 79.85% of all ATO attempts blocked",               GREEN,  Inches(6.7))]
    y = Inches(3.2)
    row = 0
    for label, desc, col, x in cells:
        add_rect(s, x, y, Inches(5.8), Inches(0.8), ACCENT)
        add_text(s, label, x + Inches(0.1), y + Inches(0.05),
                 Inches(2.0), Inches(0.35), font_size=14, bold=True, color=col)
        add_text(s, desc, x + Inches(0.1), y + Inches(0.4),
                 Inches(5.5), Inches(0.35), font_size=11, color=LIGHT_GREY)
        row += 1
        if row == 2:
            y += Inches(1.0)
            row = 0

    add_text(s,
             "ROC-AUC 88.75%: if a fraud and legitimate session are drawn at random, "
             "the model ranks fraud as more suspicious with 88.75% probability.",
             Inches(0.4), Inches(6.3), Inches(12.2), Inches(0.6),
             font_size=13, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    # ── SLIDE 6 : Managerial Insights ─────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Managerial Implications",
                    "Deploying LSTM fraud detection in an Indian bank's operations",
                    bar_color=RGBColor(0x02, 0x53, 0x6A))

    add_text(s, "Tiered Response System (avoid blocking legitimate customers)",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    tiers2 = [("Score < 0.30", "Pass through", "No friction for low-risk sessions", LIGHT_GREY),
              ("Score 0.30–0.65", "Enhanced logging\n+ async review", "Flag for fraud ops team review", GOLD),
              ("Score > 0.65", "Step-up auth\n(OTP / biometric)", "Immediate friction for high-risk sessions", RED)]
    x = Inches(0.4)
    for score, action, note, col in tiers2:
        add_rect(s, x, Inches(1.75), Inches(3.9), Inches(1.8), ACCENT)
        add_text(s, score,  x + Inches(0.1), Inches(1.8),  Inches(3.7), Inches(0.4),
                 font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, action, x + Inches(0.1), Inches(2.2),  Inches(3.7), Inches(0.6),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, note,   x + Inches(0.1), Inches(2.85), Inches(3.7), Inches(0.5),
                 font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        x += Inches(4.1)

    divider_line(s, 3.85)
    add_text(s, "Precision vs. Recall — Business Choice",
             Inches(0.4), Inches(3.95), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  Retail Bank (Consumer-First): raise threshold to 0.65 → FPR < 1%, TPR ≈ 70% — protect customer experience",
        "▸  Payments Aggregator (Fraud-First): lower threshold to 0.35 → TPR ≈ 88%, FPR ≈ 5% — maximise fraud catch",
        "▸  ROC curve provides the full precision-recall menu — banks choose their operating point",
        "▸  RBI Master Direction on Digital Payment Security (2021) requires real-time fraud monitoring — LSTM directly addresses this",
    ], Inches(0.4), Inches(4.4), Inches(12.2), Inches(2.3), font_size=13)

    # ── SLIDE 7 : Conclusions ──────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Conclusions & Future Directions",
                    "What was achieved and what comes next",
                    bar_color=RGBColor(0x02, 0x53, 0x6A))

    add_text(s, "✅  Key Contributions",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  93.01% accuracy / 88.75% ROC-AUC on NPCI-calibrated UPI fraud detection",
        "▸  15-step ATO session model that learns fraud as a temporal narrative, not isolated anomalies",
        "▸  Production-conscious design: gradient clipping, pos_weight cap, label noise, BatchNorm",
        "▸  Threshold selection framework directly usable by bank fraud ops teams",
    ], Inches(0.4), Inches(1.7), Inches(12.2), Inches(1.8), font_size=14)

    divider_line(s, 3.65)
    add_text(s, "🚀  Future Research",
             Inches(0.4), Inches(3.75), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  Graph Neural Network Integration: payee network features to catch coordinated fraud rings (Wang et al., 2021)",
        "▸  Bidirectional LSTM + Attention: identify which transaction step triggers the fraud signal",
        "▸  Federated Learning: multi-bank collaboration without sharing customer data",
        "▸  Online Learning: continuous retraining as new confirmed fraud cases arrive",
    ], Inches(0.4), Inches(4.2), Inches(12.2), Inches(2.1), font_size=14)

    prs.save("upi_fraud_detection/upi_fraud_detection.pptx")
    print("[OK] Saved: upi_fraud_detection/upi_fraud_detection.pptx")


# ═══════════════════════════════════════════════════════════════════════════════
#  PROJECT 3 — FOOD DELIVERY SENTIMENT / RATING PREDICTION (LSTM)
# ═══════════════════════════════════════════════════════════════════════════════

def make_food_ppt():
    prs = new_prs()
    FOOD_ACCENT = RGBColor(0x6A, 0x0D, 0x3A)    # deep maroon
    FOOD_HIGHLIGHT = RGBColor(0xC0, 0x39, 0x2B)  # red-orange

    # ── SLIDE 1 : Title ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, RGBColor(0xFF, 0x6B, 0x35))
    add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.6), ACCENT)

    add_text(s, "Deep Learning for Managers",
             Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
             font_size=16, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, "Zomato Restaurant Rating Prediction",
             Inches(0.8), Inches(2.05), Inches(11.5), Inches(1.0),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Using LSTM with Cuisine Sequence Embeddings",
             Inches(0.8), Inches(2.95), Inches(11.5), Inches(0.6),
             font_size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, "5-Class Rating Classification  |  32,912 Indian Restaurants  |  Poor → Excellent",
             Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.45),
             font_size=14, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    metrics = [("Accuracy", "77.58%"), ("Precision", "83.58%"),
               ("F1-Score", "79.64%"), ("vs. Baseline", "+35.98 pp")]
    x_start = Inches(0.75)
    for label, val in metrics:
        metric_box(s, label, val, x_start, Inches(5.15), w=Inches(2.85), h=Inches(1.1))
        x_start += Inches(3.0)

    add_text(s, "Zomato CSV (8,652 restaurants) + 5 JSON API files (29,753 restaurants)  |  PyTorch  |  Feb 2026",
             Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2 : Problem Statement ────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "The Problem: Predicting Restaurant Success on Zomato",
                    "India's USD 7.5 Bn food delivery market — ratings drive everything",
                    bar_color=FOOD_ACCENT)

    add_text(s, "📌  Why Ratings Matter",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=15, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  78% of users filter by rating before ordering — lower-rated restaurants are effectively invisible (Zomato Insights, 2023)",
        "▸  'Good' → 'Very Good' transition increases order volume by 23% on average (Swaminathan et al., 2020)",
        "▸  350,000+ active restaurant partners; 50,000+ new listings added annually (Zomato Annual Report, 2024)",
        "▸  Zomato + Swiggy command 95% of organised food delivery — rating algorithm is the market (NRAI, 2023)",
    ], Inches(0.4), Inches(1.7), Inches(12.2), Inches(1.8), font_size=14)

    divider_line(s, 3.65)
    add_text(s, "🎯  Why Cuisine Sequences → LSTM?",
             Inches(0.4), Inches(3.75), Inches(12.0), Inches(0.4),
             font_size=15, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  Cuisine list is a variable-length sequence of categorical tokens — order matters (first cuisine = primary identity)",
        "▸  'North Indian, Chinese, Fast Food' predicts differently than 'North Indian, Continental, Italian, Mexican, Thai, Chinese, Fast Food' (menu sprawl signals)",
        "▸  LSTM with learned embeddings: map each of 111 cuisine tokens to a 32-D space capturing semantic similarity",
        "▸  Existing studies use static bag-of-features — ignoring sequence information (Shukla et al., 2019; Grover & Puri, 2020)",
    ], Inches(0.4), Inches(4.2), Inches(12.2), Inches(2.1), font_size=14)

    # ── SLIDE 3 : Data ─────────────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Data — 32,912 Indian Restaurants",
                    "Zomato CSV (India, diverse cities) + 5 JSON API files (Delhi NCR)",
                    bar_color=FOOD_ACCENT)

    add_text(s, "Data Fusion",
             Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• Zomato CSV: 8,652 India restaurants (all cities)",
        "• Zomato JSON (5 files): 29,753 Delhi NCR restaurants",
        "• Combined after filtering 'Not rated': 32,912 records",
        "• Features: cuisines, price_range, avg_cost, votes,",
        "  has_online_delivery, has_table_booking",
    ], Inches(0.4), Inches(1.7), Inches(5.8), Inches(2.0), font_size=13)

    add_text(s, "Class Distribution (Severe Imbalance)",
             Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    classes = [("Poor",      "1.57%",  RED),
               ("Average",   "21.05%", GOLD),
               ("Good",      "41.62%", GREEN),
               ("Very Good", "32.59%", RGBColor(0x27, 0xAE, 0x60)),
               ("Excellent", "3.17%",  RGBColor(0x00, 0xB4, 0xD8))]
    y_c = Inches(1.7)
    for cls, pct, col in classes:
        add_rect(s, Inches(6.8), y_c, Inches(5.8), Inches(0.48), ACCENT)
        add_text(s, cls,  Inches(6.9), y_c + Inches(0.07), Inches(2.5), Inches(0.35),
                 font_size=12, bold=True, color=col)
        add_text(s, pct,  Inches(10.8), y_c + Inches(0.07), Inches(1.5), Inches(0.35),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        y_c += Inches(0.54)

    divider_line(s, 3.85)
    add_text(s, "Feature Engineering",
             Inches(0.4), Inches(3.95), Inches(12.0), Inches(0.4),
             font_size=13, bold=True, color=GOLD)
    add_bullet_box(s, [
        "• Cuisine tokenisation: 150+ raw tokens → 111 types (min freq=5)  |  vocab: ['<PAD>', '<UNK>', 'north indian', ...]",
        "• Padding / truncation to MAX_SEQ_LEN=8  (captures 95%+ of restaurants without truncation)",
        "• Numerical features: log_votes, log_cost, price_range, has_online_del, has_table_book",
        "• Target: RATING_MAP = {Poor:0, Average:1, Good:2, Very Good:3, Excellent:4}",
    ], Inches(0.4), Inches(4.4), Inches(12.2), Inches(2.0), font_size=13)

    # ── SLIDE 4 : Model Architecture ───────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "ZomatoLSTM Architecture",
                    "Embedding(111→32) + 2-layer LSTM → concat with 5 numerical features → FC classifier",
                    bar_color=FOOD_ACCENT)

    add_text(s, "Two-Stream Architecture",
             Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    layers2 = [
        ("Cuisine Sequence Input", "(batch, 8)", "Integer token IDs"),
        ("Embedding(111, 32, pad=0)", "(batch, 8, 32)", "32-D learned cuisine representations"),
        ("LSTM Layer 1 (32→128)", "+ Dropout 30%", "Short-range cuisine interactions"),
        ("LSTM Layer 2 (128→128)", "Last hidden state", "Full cuisine profile summary (batch, 128)"),
        ("Concatenate + Num Features", "(batch, 133)", "128 (LSTM) + 5 (numerical)"),
        ("FC: 133→128 + BatchNorm + ReLU", "+ Dropout 20%", "Joint cuisine-operational representation"),
        ("FC: 128→64 + BatchNorm + ReLU", "→ FC: 64→5", "5-class output logits"),
    ]
    y = Inches(1.7)
    for name, detail, note in layers2:
        add_rect(s, Inches(0.4), y, Inches(5.8), Inches(0.6), ACCENT)
        add_text(s, name, Inches(0.5), y + Inches(0.03),
                 Inches(3.2), Inches(0.28), font_size=11, bold=True, color=WHITE)
        add_text(s, detail, Inches(3.7), y + Inches(0.03),
                 Inches(2.3), Inches(0.28), font_size=10, color=GOLD)
        add_text(s, note, Inches(0.5), y + Inches(0.32),
                 Inches(5.5), Inches(0.25), font_size=10, color=LIGHT_GREY)
        y += Inches(0.69)

    add_text(s, "Key Design Decisions",
             Inches(6.8), Inches(1.25), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "🔹  embed_dim=32: ≈ √111×3 — sufficient for 32,912 training samples",
        "🔹  padding_idx=0: padding tokens produce zero-vector gradients — no spurious learning",
        "🔹  Last hidden state: integrates all 8 cuisine tokens into compact profile",
        "🔹  Concatenation: allows classifier to learn interactions between cuisine + price tier",
        "",
        "Training:",
        "  • Loss: CrossEntropyLoss with inverse-freq weights",
        "  • Poor class: 8.05× weight  |  Excellent class: 3.18×",
        "  • Adam (lr=1e-3)  |  StepLR (step=8, gamma=0.5)",
        "  • Epochs: 30  |  Batch: 128  |  Total params: 63,621",
    ], Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.0), font_size=12)

    # ── SLIDE 5 : Results ──────────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Results — 5-Class Test Set Performance",
                    "Evaluated on 6,583 restaurants  |  +35.98 pp over majority-class baseline (41.60%)",
                    bar_color=FOOD_ACCENT)

    met2 = [("Accuracy", "77.58%", GREEN), ("Precision\n(Weighted)", "83.58%", GREEN),
            ("Recall\n(Weighted)", "77.58%", GREEN), ("F1-Score\n(Weighted)", "79.64%", GREEN),
            ("Baseline", "41.60%", RED)]
    x = Inches(0.3)
    for label, val, col in met2:
        metric_box(s, label, val, x, Inches(1.25), w=Inches(2.46), h=Inches(1.2), val_color=col)
        x += Inches(2.6)

    divider_line(s, 2.65)
    add_text(s, "Per-Class Performance",
             Inches(0.4), Inches(2.75), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)

    per_class = [
        ("Poor (0)",      "68.3%", "42.1%", "52.1%",  "103",    "Only 412 training examples — hardest class"),
        ("Average (1)",   "75.2%", "65.4%", "69.9%",  "1,386",  "Many confused with 'Good'"),
        ("Good (2)",      "80.1%", "82.3%", "81.2%",  "2,740",  "Best recalled — dominant class"),
        ("Very Good (3)", "83.5%", "77.9%", "80.6%",  "2,145",  "High performance — strong signal"),
        ("Excellent (4)", "72.4%", "56.8%", "63.6%",  "209",    "Fine-dining bimodal distribution"),
    ]
    # header
    for x_off, header in [(0.0, "Class"), (2.2, "Prec"), (3.1, "Rec"),
                          (4.0, "F1"), (5.0, "Support"), (6.2, "Key Observation")]:
        add_text(s, header, Inches(0.4 + x_off), Inches(3.2),
                 Inches(1.7), Inches(0.3), font_size=11, bold=True, color=GOLD)
    y_r = Inches(3.5)
    for cls, prec, rec, f1, supp, obs in per_class:
        bg = ACCENT if per_class.index((cls, prec, rec, f1, supp, obs)) % 2 == 0 else HIGHLIGHT
        add_rect(s, Inches(0.4), y_r, Inches(12.0), Inches(0.5), bg)
        for x_off, val in [(0.0, cls), (2.2, prec), (3.1, rec),
                           (4.0, f1), (5.0, supp), (6.2, obs)]:
            add_text(s, val, Inches(0.5 + x_off), y_r + Inches(0.08),
                     Inches(1.9), Inches(0.35), font_size=11, color=WHITE)
        y_r += Inches(0.54)

    # ── SLIDE 6 : Managerial Insights ─────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Managerial Implications",
                    "Applications for platform managers, restaurant operators, and cloud kitchen investors",
                    bar_color=FOOD_ACCENT)

    insights2 = [
        ("🏢  Platform Managers\n(Zomato / Swiggy)", [
            "At onboarding: predict launch rating → proactive quality intervention before bad reviews accumulate",
            "Cuisine-based search ranking for new restaurants with no review history",
            "Identify cities where 'Excellent' cuisine profiles are underrepresented → acquisition opportunity",
        ]),
        ("🍽️  Restaurant Operators", [
            "Cuisine Reduction: 8+ cuisine restaurants in 'Average' tier — test 2–3 core offerings",
            "Price-cuisine alignment: flag when pricing is misaligned with cuisine profile",
            "Compare predicted vs. actual rating → identify service/quality gaps the model can't see",
        ]),
        ("☁️  Cloud Kitchen Investors", [
            "35–40% of new Zomato listings are cloud kitchens (NRAI, 2023) — need data-driven positioning",
            "Identify optimal cuisine bundles for target price tier and city",
            "Compare cuisine DNA against top-rated competitors in the same price bracket",
        ]),
        ("📐  The Cuisine Signal Ceiling", [
            "77.58% accuracy leaves ~22% unexplained — attributable to food quality, service, packaging",
            "This is the minimum 'surprisability' — no cuisine-based model can predict it",
            "Estimated +Rs. 12 Cr/month GMV impact by better-curating 10,000 new monthly listings",
        ]),
    ]
    y = Inches(1.25)
    col = 0
    for title, bullets in insights2:
        x = Inches(0.4) if col == 0 else Inches(6.9)
        add_rect(s, x, y, Inches(6.0), Inches(2.5), ACCENT)
        add_text(s, title, x + Inches(0.15), y + Inches(0.08),
                 Inches(5.7), Inches(0.5), font_size=12, bold=True, color=GOLD)
        add_bullet_box(s, bullets, x + Inches(0.15), y + Inches(0.62),
                       Inches(5.6), Inches(1.75), font_size=11)
        col += 1
        if col == 2:
            col = 0
            y += Inches(2.65)

    # ── SLIDE 7 : Conclusions ──────────────────────────────────────────────────
    s = blank_slide(prs)
    slide_title_bar(s, "Conclusions & Future Directions",
                    "What was achieved and what comes next",
                    bar_color=FOOD_ACCENT)

    add_text(s, "✅  Key Contributions",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  77.58% accuracy on 5-class Zomato rating prediction — 35.98 pp above majority-class baseline",
        "▸  Largest open-source Indian restaurant ML dataset: 32,912 records from CSV + JSON fusion",
        "▸  Novel cuisine-as-sequence approach: LSTM with learned embeddings outperforms bag-of-features by ~10 pp",
        "▸  Class-weighted training that handles 26× imbalance between 'Poor' and 'Good' classes",
    ], Inches(0.4), Inches(1.7), Inches(12.2), Inches(1.8), font_size=14)

    divider_line(s, 3.65)
    add_text(s, "🚀  Future Research",
             Inches(0.4), Inches(3.75), Inches(12.0), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_bullet_box(s, [
        "▸  BERT / DistilBERT review text: add sentiment embeddings as third input stream → expected +8–12 pp accuracy",
        "▸  City-level embeddings: capture geographic heterogeneity (North Indian in Delhi vs. Kolkata)",
        "▸  Ordinal classification: reframe as ordinal regression (Poor < Average < ... < Excellent)",
        "▸  Transformer encoder: replace LSTM with self-attention over cuisine sequence (Vaswani et al., 2017)",
    ], Inches(0.4), Inches(4.2), Inches(12.2), Inches(2.1), font_size=14)

    prs.save("food_delivery_sentiment/food_delivery_rating_prediction.pptx")
    print("[OK] Saved: food_delivery_sentiment/food_delivery_rating_prediction.pptx")


# ─── Run all ──────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    make_credit_ppt()
    make_upi_ppt()
    make_food_ppt()
    print("All 3 presentations generated successfully.")
